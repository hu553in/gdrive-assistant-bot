from __future__ import annotations

import io
import subprocess
from types import SimpleNamespace

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from gdrive_assistant_bot.extractors.office.excel import XlsExtractor, XlsxExtractor
from gdrive_assistant_bot.extractors.office.powerpoint import PptExtractor, PptxExtractor
from gdrive_assistant_bot.extractors.office.word import DocExtractor, DocxExtractor
from gdrive_assistant_bot.extractors.pdf import PDFExtractor
from gdrive_assistant_bot.extractors.text import TextBasedFileExtractor

_INT_3 = 3
_INT_4 = 4
_INT_5 = 5
_INT_7 = 7
_INT_8 = 8
_INT_9 = 9
_INT_11 = 11
_INT_12 = 12


def test_text_extractor_properties_and_helpers() -> None:
    extractor = TextBasedFileExtractor()
    assert "application/json" in extractor.mime_types
    assert "py" in extractor.file_extensions
    assert extractor.mime_prefixes == ["text/"]
    assert extractor._to_int(_INT_5) == _INT_5
    assert extractor._to_int(str(_INT_7)) == _INT_7
    assert extractor._to_int("x") is None
    assert extractor._extension({"name": "README.MD"}) == "md"
    assert extractor._extension({"name": "noext"}) is None
    assert extractor._normalized_file_type("js") == "javascript"
    assert extractor._normalized_file_type("ts") == "typescript"
    assert extractor._normalized_file_type("yaml") == "yaml"
    assert extractor._normalized_file_type("md") == "markdown"
    assert extractor._normalized_file_type("json") == "json"
    assert extractor._normalized_file_type("toml") == "toml"
    assert extractor._normalized_file_type("sh") == "shell"
    assert extractor._normalized_file_type("csv") == "csv"
    assert extractor._normalized_file_type("bin") == "text"


def test_pdf_extract_text_pypdf_with_limit(monkeypatch: pytest.MonkeyPatch) -> None:
    class _Page:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class _Reader:
        def __init__(self, _stream: io.BytesIO) -> None:
            self.pages = [_Page("first"), _Page("second")]

    monkeypatch.setattr("gdrive_assistant_bot.extractors.pdf.PdfReader", _Reader)
    text = PDFExtractor._extract_text_pypdf(b"pdf", max_pages=1)
    assert "first" in text
    assert "limited to 1 pages" in text


def test_pdf_extract_text_pdfplumber_with_limit(monkeypatch: pytest.MonkeyPatch) -> None:
    class _Page:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class _PdfCtx:
        def __init__(self) -> None:
            self.pages = [_Page("one"), _Page("two")]

        def __enter__(self) -> _PdfCtx:
            return self

        def __exit__(self, *_args) -> None:
            return None

    fake_module = SimpleNamespace(open=lambda _stream: _PdfCtx())
    monkeypatch.setattr("gdrive_assistant_bot.extractors.pdf.pdfplumber", fake_module)
    text = PDFExtractor._extract_text_pdfplumber(b"pdf", max_pages=1)
    assert "one" in text
    assert "limited to 1 pages" in text


def test_pdf_extract_text_from_pdf_invalid_engine() -> None:
    extractor = PDFExtractor()
    with pytest.raises(ValueError, match="Unsupported PDF extraction engine"):
        extractor._extract_text_from_pdf(b"pdf", max_pages=1, engine="unknown")


def test_word_docx_extracts_paragraphs_and_tables() -> None:
    doc = Document()
    doc.add_paragraph("Paragraph")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    buffer = io.BytesIO()
    doc.save(buffer)

    text = DocxExtractor._extract_docx(buffer.getvalue())
    assert "Paragraph" in text
    assert "A | B" in text
    assert DocxExtractor().mime_types
    assert DocxExtractor().file_extensions
    assert DocxExtractor._extension({"name": "notes.docx"}) == "docx"
    assert DocxExtractor._to_int(str(_INT_11)) == _INT_11


def test_word_doc_raises_on_catdoc_failure(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(
        subprocess,
        "run",
        lambda *_args, **_kwargs: SimpleNamespace(returncode=1, stdout=b"", stderr=b"broken"),
    )
    with pytest.raises(RuntimeError, match="catdoc failed"):
        DocExtractor._extract_doc(b"binary")
    assert DocExtractor().mime_types
    assert DocExtractor().file_extensions
    assert DocExtractor._extension({"name": "legacy.doc"}) == "doc"
    assert DocExtractor._to_int(str(_INT_12)) == _INT_12


def test_xlsx_parser_extracts_rows_and_limits() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Data"
    sheet.append(["A", "B"])
    sheet.append(["1", "2"])
    second = workbook.create_sheet("Second")
    second.append(["x"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    text = XlsxExtractor._extract_xlsx(buffer.getvalue(), max_sheets=1, max_rows=1)
    assert "=== SHEET: Data ===" in text
    assert "A\tB" in text
    assert "limited to 1 rows" in text
    assert "limited to 1 sheets" in text
    assert XlsxExtractor().mime_types
    assert XlsxExtractor().file_extensions
    assert XlsxExtractor._extension({"name": "book.xlsx"}) == "xlsx"
    assert XlsxExtractor._to_int(str(_INT_9)) == _INT_9


def test_xls_parser_extracts_rows_and_limits(monkeypatch: pytest.MonkeyPatch) -> None:
    class _Sheet:
        name = "Legacy"
        nrows = 3
        ncols = 2

        def cell_value(self, row: int, col: int) -> str:
            data = {(0, 0): "A", (0, 1): "B", (1, 0): "1", (1, 1): "2", (2, 0): "3", (2, 1): "4"}
            return data[(row, col)]

    class _Workbook:
        nsheets = 2

        def sheet_by_index(self, _idx: int) -> _Sheet:
            return _Sheet()

    monkeypatch.setattr("xlrd.open_workbook", lambda **_kwargs: _Workbook())
    text = XlsExtractor._extract_xls(b"xls", max_sheets=1, max_rows=2)
    assert "=== SHEET: Legacy ===" in text
    assert "A\tB" in text
    assert "1\t2" in text
    assert "limited to 2 rows" in text
    assert "limited to 1 sheets" in text
    assert XlsExtractor().mime_types
    assert XlsExtractor().file_extensions
    assert XlsExtractor._extension({"name": "book.xls"}) == "xls"
    assert XlsExtractor._to_int(str(_INT_8)) == _INT_8


def test_pptx_parser_extracts_shapes_and_tables() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "Headline"
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    buffer = io.BytesIO()
    presentation.save(buffer)

    text = PptxExtractor._extract_pptx(buffer.getvalue())
    assert "=== SLIDE 1 ===" in text
    assert "Headline" in text
    assert "A | B" in text
    assert PptxExtractor().mime_types
    assert PptxExtractor().file_extensions
    assert PptxExtractor._extension({"name": "deck.pptx"}) == "pptx"
    assert PptxExtractor._to_int(str(_INT_3)) == _INT_3


def test_ppt_parser_raises_on_catppt_failure(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(
        subprocess,
        "run",
        lambda *_args, **_kwargs: SimpleNamespace(returncode=1, stdout=b"", stderr=b"broken"),
    )
    with pytest.raises(RuntimeError, match="catppt failed"):
        PptExtractor._extract_ppt(b"binary")
    assert PptExtractor().mime_types
    assert PptExtractor().file_extensions
    assert PptExtractor._extension({"name": "legacy.ppt"}) == "ppt"
    assert PptExtractor._to_int(str(_INT_4)) == _INT_4
