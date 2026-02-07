from __future__ import annotations

import io
import os
import subprocess
import tempfile
from typing import Any

from pptx import Presentation

from ..base import ExtractedContent, ExtractionContext, FileExtractor

PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PPT_MIME_TYPE = "application/vnd.ms-powerpoint"


class PptxExtractor(FileExtractor):
    """Extract text from PPTX files."""

    @property
    def mime_types(self) -> list[str]:
        return [PPTX_MIME_TYPE]

    @property
    def file_extensions(self) -> list[str]:
        return ["pptx"]

    def can_extract(self, file_meta: dict[str, Any]) -> bool:
        if file_meta.get("mimeType") == PPTX_MIME_TYPE:
            return True
        return self._extension(file_meta) == "pptx"

    def extract(self, file_meta: dict[str, Any], context: ExtractionContext) -> ExtractedContent:
        size = self._to_int(file_meta.get("size"))
        max_bytes = int(context.settings.OFFICE_MAX_FILE_SIZE_MB * 1024 * 1024)
        if size and size > max_bytes:
            return ExtractedContent(
                text="", file_type="pptx", metadata={"skipped": "size_limit", "size_bytes": size}
            )

        pptx_bytes = context.download_binary(file_meta["id"])
        text = self._extract_pptx(pptx_bytes)
        return ExtractedContent(
            text=text.strip(),
            file_type="pptx",
            metadata={"mime_type": file_meta.get("mimeType"), "file_size_bytes": len(pptx_bytes)},
        )

    @staticmethod
    def _extract_pptx(pptx_bytes: bytes) -> str:
        presentation = Presentation(io.BytesIO(pptx_bytes))
        lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"=== SLIDE {slide_index} ===")
            for shape in slide.shapes:
                lines.extend(PptxExtractor._extract_shape_lines(shape))
            lines.append("")
        return "\n".join(lines)

    @staticmethod
    def _extract_shape_lines(shape: Any) -> list[str]:
        lines: list[str] = []

        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                lines.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))

        child_shapes = getattr(shape, "shapes", None)
        if child_shapes:
            for child in child_shapes:
                lines.extend(PptxExtractor._extract_shape_lines(child))

        return lines

    @staticmethod
    def _extension(file_meta: dict[str, Any]) -> str | None:
        ext = file_meta.get("fileExtension")
        if isinstance(ext, str) and ext.strip():
            return ext.lower().lstrip(".")

        name = file_meta.get("name")
        if not isinstance(name, str) or "." not in name:
            return None
        return name.rsplit(".", 1)[-1].strip().lower() or None

    @staticmethod
    def _to_int(value: Any) -> int | None:
        if isinstance(value, int):
            return value
        if isinstance(value, str) and value.isdigit():
            return int(value)
        return None


class PptExtractor(FileExtractor):
    """Extract text from legacy PPT files."""

    @property
    def mime_types(self) -> list[str]:
        return [_PPT_MIME_TYPE]

    @property
    def file_extensions(self) -> list[str]:
        return ["ppt"]

    def can_extract(self, file_meta: dict[str, Any]) -> bool:
        if file_meta.get("mimeType") == _PPT_MIME_TYPE:
            return True
        return self._extension(file_meta) == "ppt"

    def extract(self, file_meta: dict[str, Any], context: ExtractionContext) -> ExtractedContent:
        size = self._to_int(file_meta.get("size"))
        max_bytes = int(context.settings.OFFICE_MAX_FILE_SIZE_MB * 1024 * 1024)
        if size and size > max_bytes:
            return ExtractedContent(
                text="", file_type="ppt", metadata={"skipped": "size_limit", "size_bytes": size}
            )

        ppt_bytes = context.download_binary(file_meta["id"])
        text = self._extract_ppt(ppt_bytes)
        return ExtractedContent(
            text=text.strip(),
            file_type="ppt",
            metadata={"mime_type": file_meta.get("mimeType"), "file_size_bytes": len(ppt_bytes)},
        )

    @staticmethod
    def _extract_ppt(ppt_bytes: bytes) -> str:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
            tmp.write(ppt_bytes)
            tmp_path = tmp.name

        try:
            result = subprocess.run(["catppt", tmp_path], capture_output=True, check=False)
        except FileNotFoundError as exc:
            raise RuntimeError("Legacy PPT extraction requires the 'catppt' binary.") from exc
        finally:
            os.unlink(tmp_path)

        if result.returncode != 0:
            stderr = (result.stderr or b"").decode("utf-8", errors="replace").strip()
            raise RuntimeError(f"catppt failed to extract PPT file: {stderr or 'unknown error'}")

        return (result.stdout or b"").decode("utf-8", errors="replace")

    @staticmethod
    def _extension(file_meta: dict[str, Any]) -> str | None:
        ext = file_meta.get("fileExtension")
        if isinstance(ext, str) and ext.strip():
            return ext.lower().lstrip(".")

        name = file_meta.get("name")
        if not isinstance(name, str) or "." not in name:
            return None
        return name.rsplit(".", 1)[-1].strip().lower() or None

    @staticmethod
    def _to_int(value: Any) -> int | None:
        if isinstance(value, int):
            return value
        if isinstance(value, str) and value.isdigit():
            return int(value)
        return None
